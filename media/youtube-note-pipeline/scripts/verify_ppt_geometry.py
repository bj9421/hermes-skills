#!/usr/bin/env python3
"""PPT 幾何驗證腳本（標準化）"""
import sys, glob, os
sys.path.insert(0, '/opt/data/skills/media/youtube-note-pipeline/scripts')
from pptx import Presentation
from pptx.util import Emu

def verify_ppt_geometry(pptx_path: str) -> int:
    """驗證 PPT 所有文字框不超界，回傳問題數量。"""
    prs = Presentation(pptx_path)
    SW, SH = prs.slide_width, prs.slide_height
    issues = 0

    def box(s):
        return (s.left, s.top, s.left + s.width, s.top + s.height)

    for i, slide in enumerate(prs.slides, 1):
        shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
        for s in shapes:
            l, t, r, b = box(s)
            tol = 45720  # 0.05 inch tolerance
            if l < -tol or t < -tol or r > SW + tol or b > SH + tol:
                print(f'  ❌ Slide {i} 超界: {s.text_frame.text.strip()[:20]!r}')
                issues += 1
        # 文字框重疊檢查
        for a_i in range(len(shapes)):
            for b_i in range(a_i + 1, len(shapes)):
                a, b = shapes[a_i], shapes[b_i]
                la, ta, ra, ba = box(a)
                lb, tb, rb, bb = box(b)
                overlap_w = min(ra, rb) - max(la, lb)
                overlap_h = min(ba, bb) - max(ta, tb)
                if overlap_w > 0 and overlap_h > 0:
                    ta_, tb_ = a.text_frame.text.strip()[:15], b.text_frame.text.strip()[:15]
                    print(f'  ⚠️ Slide {i} 文字框重疊: {ta_!r} × {tb_!r}')
    return issues

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('用法: python verify_ppt_geometry.py <path.pptx>')
        sys.exit(1)
    path = sys.argv[1]
    issues = verify_ppt_geometry(path)
    print(f'超界問題: {issues}')
    if issues == 0:
        print('✅ 幾何驗證完成')
    else:
        print(f'❌ {issues} 個超界')
        sys.exit(1)