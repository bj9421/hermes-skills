#!/usr/bin/env python3
"""Render PPTX slides to PNG previews (no LibreOffice needed).

WHY: Telegram cannot preview .pptx attachments — it sends them as a bare file
with no visual thumbnail, so the user reports "ppt 還是沒看到" even though the
file exists and is valid. Send slide PNGs alongside the .pptx so the user can
actually see the content.

NOTE: This is a CONTENT preview, not a faithful layout render. It redraws the
text from every shape on a dark background (960x540 px). Colors, images,
shapes and precise layout are NOT reproduced — good enough to confirm content
exists and is readable, not a substitute for PowerPoint/LibreOffice rendering.

Usage:
    python ppt_preview_render.py <path.pptx> [out_dir]
    # out_dir defaults to ./ppt_preview; outputs slide_01.png ... slide_NN.png
"""
import os
import re
import sys
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

FONT_CANDIDATES = [
    "/opt/data/fonts/SourceHanSansTC-Regular.otf",
    "/opt/data/fonts/NotoSansCJKtc-Regular.otf",
    "/opt/data/fonts/NotoSansTC-Regular.otf",
    "/opt/data/fonts/NotoSansSC-Bold.ttf",
]
EMOJI_FONT_PATH = "/opt/data/fonts/NotoEmoji-Regular.ttf"

BG = (13, 24, 44)          # deep navy (matches notehub dark scheme)
FG = (255, 255, 255)       # white
ACCENT = (255, 180, 90)    # orange for bold/headers

# emoji / 符號 range：讓預覽圖用 NotoEmoji 字型渲染，避免中文字型無 glyph 顯示成方框
EMOJI_RE = re.compile(r"[\U0001F000-\U0001FAFF\u2600-\u27BF\u2B00-\u2BFF\uFE0F\u2460-\u2473]")


def load_font(size: int):
    for p in FONT_CANDIDATES:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except Exception:
                continue
    return ImageFont.load_default()


def load_emoji_font(size: int):
    try:
        return ImageFont.truetype(EMOJI_FONT_PATH, size)
    except Exception:
        return None


def draw_mixed(draw, xy, text, font, emoji_font, fill):
    """混合字型渲染：emoji 字元用 NotoEmoji，其餘用主字型。回傳結束 x 座標。"""
    x, y = xy
    if not emoji_font:
        draw.text((x, y), text, font=font, fill=fill)
        return x + draw.textlength(text, font=font)
    segments = []
    cur, cur_is_emoji = "", None
    for ch in text:
        is_emoji = bool(EMOJI_RE.match(ch))
        if cur_is_emoji is None:
            cur_is_emoji = is_emoji
        if is_emoji != cur_is_emoji:
            segments.append((cur, cur_is_emoji))
            cur, cur_is_emoji = ch, is_emoji
        else:
            cur += ch
    segments.append((cur, cur_is_emoji))
    for seg, is_emoji in segments:
        f = emoji_font if is_emoji else font
        draw.text((x, y), seg, font=f, fill=fill)
        x += draw.textlength(seg, font=f)
    return x


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 1
    pptx_path = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else "ppt_preview"
    os.makedirs(out_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    w = int(prs.slide_width / 914400 * 96)   # EMU -> px (96 dpi)
    h = int(prs.slide_height / 914400 * 96)
    print(f"Slides: {len(prs.slides)}, size: {w}x{h}", file=sys.stderr)

    for i, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (w, h), BG)
        draw = ImageDraw.Draw(img)
        texts = []  # (text, size_pt, bold)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs)
                if not txt.strip():
                    continue
                size = None
                bold = False
                for r in para.runs:
                    if r.font.size:
                        size = r.font.size.pt
                    bold = bold or bool(r.font.bold)
                texts.append((txt.strip(), size or 18, bold))

        y = 60
        for txt, size, bold in texts:
            font = load_font(int(size * 1.4))
            emoji_font = load_emoji_font(int(size * 1.4))
            # manual wrap to fit width
            lines, cur = [], ""
            for ch in txt:
                if draw.textlength(cur + ch, font=font) > w - 160:
                    lines.append(cur)
                    cur = ch
                else:
                    cur += ch
            lines.append(cur)
            color = ACCENT if bold else FG
            for ln in lines:
                draw_mixed(draw, (80, y), ln, font, emoji_font, color)
                y += int(size * 1.4) + 8

        out = os.path.join(out_dir, f"slide_{i:02d}.png")
        img.save(out)
        print(f"  Slide {i}: {out}", file=sys.stderr)

    print(f"DONE: {os.path.abspath(out_dir)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
