#!/usr/bin/env python3
"""
ppt_gen.py — Generate PowerPoint slides from podcast scripts.

Extracts key points via LLM and creates a professional presentation.
"""

import os
import sys
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# JSON 容錯解析（2026-08-06 階段 1：LLM 輸出格式不穩，直接 json.loads 常失敗
# → fallback「無法提取」→ 簡報品質爛。強化解析，成功率高很多。）
# ---------------------------------------------------------------------------
def _parse_json_loose(text: str) -> dict | None:
    """容錯解析 LLM 輸出的 JSON：跳過前綴雜訊、找到第一個完整 JSON 物件。

    策略：
    1. strip + 找第一個 '{' 開始
    2. raw_decode 嘗試完整解析（可跳過尾部雜訊）
    3. 失敗 → 找最後一個 '}' 截斷再試（LLM 常在結尾附加說明文字）
    """
    if not text:
        return None
    import json as _json
    t = text.strip()
    start = t.find('{')
    if start == -1:
        return None
    t = t[start:]
    try:
        obj, _ = _json.JSONDecoder().raw_decode(t)
        return obj
    except _json.JSONDecodeError:
        end = t.rfind('}')
        if end == -1:
            return None
        try:
            return _json.loads(t[:end + 1])
        except _json.JSONDecodeError:
            return None


# ---------------------------------------------------------------------------
# Extract key points from script via LLM
# ---------------------------------------------------------------------------
def _extract_key_points(script: str, title: str, lang: str = "zh") -> dict:
    """Extract structured key points from podcast script.
    
    Returns dict with:
        - title: presentation title
        - subtitle: subtitle/tagline
        - points: list of {heading, bullets} dicts
        - summary: conclusion paragraph
    """
    # ⚠️ 2026-07-31 使用者指示：LLM 整理文檔一律用 Zen，不用 NVIDIA（NVIDIA 僅供 Whisper）
    try:
        # 🔴 2026-08-06：改用 call_llm（Zen→AGNES→Groq fallback 鏈）。
        # 原本 call_zen 無 fallback — Zen timeout 時直接變 3 slides 基本版。
        from notehub.core.llm import call_llm
    except ImportError:
        call_llm = None
    if not call_llm:
        print("[WARN] call_llm unavailable — cannot extract key points", file=sys.stderr)
        return None

    lang_hint = "使用繁體中文" if lang in ("zh", "zh-TW") else "Use English"

    # 🔴 2026-08-06 質量升級（階段 1）：
    # 採用 104 職場力「投影片大綱規劃模組」精神 + 2Slides 專業技巧 + one-idea-per-slide 原則。
    # 保留既有 JSON 結構（title/subtitle/points/summary）→ 渲染端零改動，質量規則大幅提升。
    prompt = f"""從以下口播腳本中提取重點，產生簡報用的結構化資料。

{lang_hint}

【角色】你是簡報架構師（Presentation Architect），擅長把長篇內容濃縮成有故事線、每頁一個核心訊息的專業簡報。

【內部思考流程】（只思考，不要輸出思考過程）
1. 找出整份內容的敘事弧線：Hook（吸睛開場）→ 問題/好奇 → 關鍵論點/證據 → 案例/數據 → 行動/總結
2. 每一頁只傳達「一個核心訊息」（one-idea-per-slide），不要塞多個概念
3. 標題採用 tagline 形式：動詞開頭或問句，具體不空泛

輸出嚴格 JSON 格式（不要加 markdown code fence）：
{{
  "title": "簡報標題（簡潔有力，10 字內最佳）",
  "subtitle": "副標題或一句話摘要（點出核心價值）",
  "points": [
    {{
      "heading": "重點標題（tagline 式，動詞/問句，≤10 字）",
      "slide_type": "版型類型（見下方規則）",
      "bullets": ["要點1", "要點2"]
    }}
  ],
  "summary": "結論段落（2-3句話，含一個明確 takeaway 與行動建議）"
}}

【slide_type 版型規則】（每頁依內容本質選最合適的一種，不要每頁都 content）
- "hook"：開場吸睛頁（第 1 頁）：一個震撼事實、矛盾或提問當 heading，bullets 可留空或 1 條補充
- "content"：一般內文頁：2-4 條 bullets，每條一個要點
- "data"：數據頁：bullets 每條 = 「具體數字 + 一句說明」（如「70% 使用者一週內回訪」），數字帶單位（%、倍、萬…）
- "quote"：金句頁：heading 就是一句有力金句（≤20 字），bullets 可留空或 1 條說明出處
- "qa"：QA 頁：bullets = 3 個建議問題（問題形式，❓不寫），最後 1 條是 CTA 行動呼籲
- "action"：行動/總結頁：bullets = 具體行動步驟或下一步建議
- "comparison"：對比頁：bullets 分兩組（左邊 vs 右邊，用「→」分隔），如「舊方法 vs 新方法」
- "timeline"：時間軸/流程頁：bullets = 步驟順序（每條一個時間點或步驟），可含「→」前綴
- "split"：左右分頁：heading = 左區標題，bullets[0] = 右區文字（案例說明、詳細描述）

規則：
- points 產出 5-7 個重點，依敘事弧線排列：第 1 個是 hook，中間是 content/data/quote/comparison/timeline 混搭，最後 1 個是 action
- 5-7 頁中至少包含：1 個 data、1 個 quote、1 個 action（qa 視內容需要）
- 當內容涉及「兩個選項比較」時用 comparison（如 A vs B、舊版 vs 新版）
- 當內容涉及「時間順序/流程/步驟」時用 timeline
- 每個 heading 只代表一個核心訊息，彼此不重複
- content/data 的 bullets 每條不超過 15 字，要具體有畫面（用數字、名稱、對比），避免流水帳與形容詞堆疊
- quote 的 heading 是金句本身（可 12-20 字），不套用 ≤10 字規則
- summary 總結核心 takeaway，並給一句行動建議（「你可以…」「下一步…」）
- 不要臆測腳本沒提到的內容
- 🔴 重要：所有內容文字禁止使用 ASCII 雙引號（"），一律用中文引號「」或直接不用，避免破壞 JSON

口播腳本：
{script[:6000]}
"""
    try:
        result = call_llm(
            [
                {"role": "system", "content": "你是結構化資料提取專家，只輸出 JSON。"},
                {"role": "user", "content": prompt},
            ],
            # 🔴 2026-08-06：max_tokens 設 0（不帶）— deepseek-v4-flash 是 reasoning
            # 模型，設 max_tokens 會被思考過程吃光 → content 空 → 回 None（同 podcast.py 做法）
            max_tokens=0,
            temperature=0.3,
        )
        if result:
            # 🔴 2026-08-06 階段 1：改用容錯解析（raw_decode + rfind fallback），
            #   取代原本只 strip code fence + json.loads（LLM 格式瑕疵就整段失敗）。
            data = _parse_json_loose(result)
            if data:
                # 🔴 2026-08-06 階段 1：部分 model 回 key_points 結構（title/description）
                #   → normalize 成渲染端吃的 points（heading/bullets）結構。
                if 'points' not in data and isinstance(data.get('key_points'), list):
                    data['points'] = []
                    for kp in data['key_points']:
                        bullets = []
                        if kp.get('description'):
                            bullets.append(kp['description'])
                        if kp.get('details') and isinstance(kp['details'], list):
                            bullets.extend(str(d) for d in kp['details'])
                        data['points'].append({
                            'heading': kp.get('title', ''),
                            'slide_type': kp.get('slide_type', 'content'),
                            'bullets': bullets,
                        })
                # 🔴 2026-08-06 階段 2：沒有 slide_type 的舊結構 → 預設 content
                for pt in data.get('points', []):
                    if not pt.get('slide_type'):
                        pt['slide_type'] = 'content'
                return data
    except Exception as e:
        print(f"[WARN] Key point extraction failed: {e}", file=sys.stderr)
    return None


# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
_DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Dark navy
_ACCENT = RGBColor(0xE8, 0x4D, 0x3D)         # Warm red
_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
_CARD_BG = RGBColor(0x2D, 0x2D, 0x44)        # Slightly lighter navy

# 🔴 2026-08-07：配色方案選擇
COLOR_SCHEMES = {
    "dark": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0xE8, 0x4D, 0x3D),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_light": RGBColor(0xCC, 0xCC, 0xCC),
        "card_bg": RGBColor(0x2D, 0x2D, 0x44),
    },
    "blue": {
        "bg": RGBColor(0x0F, 0x1C, 0x3F),
        "accent": RGBColor(0x4D, 0x9B, 0xE8),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_light": RGBColor(0xB8, 0xC5, 0xD6),
        "card_bg": RGBColor(0x1A, 0x2A, 0x4A),
    },
    "green": {
        "bg": RGBColor(0x0D, 0x2B, 0x1E),
        "accent": RGBColor(0x4D, 0xC8, 0x78),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "text_light": RGBColor(0xB8, 0xD6, 0xC5),
        "card_bg": RGBColor(0x1A, 0x3C, 0x2E),
    },
    "light": {
        "bg": RGBColor(0xF5, 0xF5, 0xF5),
        "accent": RGBColor(0xE8, 0x4D, 0x3D),
        "text_white": RGBColor(0x1A, 0x1A, 0x2E),
        "text_light": RGBColor(0x66, 0x66, 0x66),
        "card_bg": RGBColor(0xE8, 0xE8, 0xE8),
    },
}

# 🔴 2026-08-06：繁中字型 — 優先思源黑體（Adobe/Google，SIL OFL 免費商用），備選 Noto Sans CJK TC
_FONT_NAME = "Source Han Sans TC"


def _apply_cjk_font(prs: Presentation, name: str = _FONT_NAME):
    """遍歷全部 slide 文字 run，設定 latin + East Asian 字型。

    ⚠️ python-pptx 的 run.font.name 只設 latin typeface，中文需要額外設 a:ea
    （East Asian）屬性才生效 — 否則 PowerPoint 用系統預設中文字型（不可控）。
    """
    from pptx.oxml.ns import qn
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = name  # latin
                    rPr = run._r.get_or_add_rPr()
                    ea = rPr.find(qn('a:ea'))
                    if ea is None:
                        ea = rPr.makeelement(qn('a:ea'), {})
                        rPr.append(ea)
                    ea.set('typeface', name)



# ---------------------------------------------------------------------------
# Emoji Enhancement (2026-08-07)
# ---------------------------------------------------------------------------
_EMOJI_FONT = "Noto Color Emoji"
_EMOJI_KEYWORD_MAP = {
    # 🩺 健康醫療（2026-08-07 新增：K2 案例發現原表只涵蓋科技/AI 主題）
    "維生素": "💊", "维他命": "💊", "补充剂": "💊", "补充": "💊", "保健": "💊", "vitamin": "💊", "supplement": "💊",
    "鈣質": "🦴", "钙质": "🦴", "骨骼": "🦴", "骨質": "🦴", "骨密度": "🦴", "bone": "🦴", "calcium": "🦴",
    "心臟": "❤️", "心脏": "❤️", "心血管": "❤️", "血管": "❤️", "冠狀動脈": "❤️", "heart": "❤️", "cardiovascular": "❤️",
    "血液": "🩸", "血压": "🩸", "血壓": "🩸", "blood": "🩸",
    "健康": "🌿", "養生": "🌿", "养生": "🌿", "health": "🌿", "wellness": "🌿",
    "試驗": "🔬", "试验": "🔬", "研究": "🔬", "临床": "🔬", "臨床": "🔬", "trial": "🔬", "research": "🔬", "study": "🔬",
    "劑量": "⚖️", "剂量": "⚖️", "dose": "⚖️", "dosage": "⚖️",
    "医生": "🩺", "醫生": "🩺", "医师": "🩺", "醫師": "🩺", "医療": "🩺", "医疗": "🩺", "doctor": "🩺",
    "ai": "🤖", "artificial intelligence": "🤖", "机器学习": "🤖", "深度学习": "🤖",
    "data": "📊", "数据": "📊", "statistics": "📊", "analytics": "📊",
    "idea": "💡", "创新": "💡", "insight": "💡", "关键": "💡",
    "rocket": "🚀", "增长": "🚀", "效率": "🚀", "提升": "🚀",
    "check": "✅", "成功": "✅", "验证": "✅", "正确": "✅",
    "question": "❓", "疑问": "❓", "问题": "❓",
    "lightning": "⚡", "快速": "⚡", "速度": "⚡", "性能": "⚡",
    "fire": "🔥", "热门": "🔥", "趋势": "🔥", "火爆": "🔥",
    "star": "🌟", "明星": "🌟", "最佳": "🌟", "顶级": "🌟",
    "chat": "💬", "对话": "💬", "讨论": "💬", "交流": "💬",
    "target": "🎯", "目标": "🎯", "核心": "🎯", "焦点": "🎯",
    "chart": "📈", "上升": "📈", "下降": "📉",
    "gears": "⚙️", "系统": "⚙️", "技术": "⚙️", "架构": "⚙️",
    "book": "📚", "学习": "📚", "知识": "📚", "教程": "📚",
    "video": "🎬", "视频": "🎬", "教程": "🎬",
    "phone": "📱", "移动": "📱", "应用": "📱",
    "cloud": "☁️", "云": "☁️", "计算": "☁️",
    "database": "🗄️", "存储": "🗄️", "数据库": "🗄️",
    "security": "🔒", "安全": "🔒", "保护": "🔒",
    "eye": "👁️", "视觉": "👁️", "图像": "👁️",
    "brain": "🧠", "智能": "🧠", "思考": "🧠",
    "robot": "🤖", "自动化": "🤖",
    "warning": "⚠️", "注意": "⚠️", "警告": "⚠️",
    "time": "⏰", "时间": "⏰", "流程": "⏰",
    "compare": "⚖️", "对比": "⚖️", "比较": "⚖️",
    "link": "🔗", "链接": "🔗", "连接": "🔗",
    "money": "💰", "成本": "💰", "价格": "💰",
    "people": "👥", "用户": "👥", "团队": "👥",
    "settings": "🔧", "设置": "🔧", "配置": "🔧",
    "gift": "🎁", "免费": "🎁", "优惠": "🎁",
    "bell": "🔔", "通知": "🔔", "提醒": "🔔",
    "envelope": "📧", "邮件": "📧", "发送": "📧",
    "camera": "📷", "图片": "📷", "照片": "📷",
    "microphone": "🎤", "音频": "🎤", "声音": "🎤",
    "tv": "📺", "电视": "📺", "播放": "📺",
    "computer": "💻", "电脑": "💻", "桌面": "💻",
    "game": "🎮", "游戏": "🎮", "娱乐": "🎮",
    "music": "🎵", "音乐": "🎵", "音频": "🎵",
    "film": "🎞️", "电影": "🎞️", "影片": "🎞️",
    "newspaper": "📰", "新闻": "📰", "报道": "📰",
    "calendar": "📅", "日期": "📅", "时间": "📅",
    "location": "📍", "位置": "📍", "地图": "📍",
    "home": "🏠", "首页": "🏠", "主页": "🏠",
    "shopping": "🛒", "购物": "🛒", "购买": "🛒",
    "heart": "❤️", "喜欢": "❤️", "收藏": "❤️",
    "thumbsup": "👍", "推荐": "👍", "点赞": "👍",
    "thumbsdown": "👎", "反对": "👎",
    "wave": "👋", "欢迎": "👋", "你好": "👋",
    "peace": "✌️", "胜利": "✌️",
    "vulcan": "🖖", "科学": "🖖",
    "raised_hand": "✋", "停止": "✋",
    "ok_hand": "👌", "完美": "👌",
    "point_up": "☝️", "重点": "☝️",
    "point_down": "👇", "下方": "👇",
    "crossed_fingers": "🤞", "幸运": "🤞",
    "love": "💕", "爱": "💕",
    "anger": "💢", "愤怒": "💢",
    "bomb": "💣", "爆炸": "💣",
    "sweat": "💦", "出汗": "💦",
    " dizzy": "😵", "晕": "😵",
    "smile": "😄", "快乐": "😄",
    "sob": "😭", "感动": "😭",
    "sunglasses": "😎", "酷": "😎",
    "nerd": "🤓", "极客": "🤓",
    "smirk": "😏", "微笑": "😏",
    "thinking": "🤔", "思考": "🤔",
    "nauseated": "🤢", "恶心": "🤢",
    "sneezing": "🤧", "感冒": "🤧",
    "hot_face": "🥵", "热": "🥵",
    "cold_face": "🥶", "冷": "🥶",
    "angry": "😠", "生气": "😠",
    "rage": "😡", "愤怒": "😡",
    "cry": "😢", "哭泣": "😢",
    "disappointed": "😞", "失望": "😞",
    "worried": "😟", "担心": "😟",
    "yum": "😋", "好吃": "😋",
    "relieved": "😌", "放松": "😌",
    "innocent": "😇", "天真": "😇",
    "heart_eyes": "😍", "喜爱": "😍",
    "kissing": "😘", "亲吻": "😘",
    "wink": "😉", "眨眼": "😉",
    "relaxed": "😌", "放松": "😌",
    "laughing": "😆", "大笑": "😆",
    "sweat_smile": "😅", "汗": "😅",
    "cool": "😎", "酷": "😎",
    "smiley": "😃", "开心": "😃",
    "grin": "😁", "咧嘴": "😁",
    "hugging": "🤗", "拥抱": "🤗",
    "hushed": "🤫", "安静": "🤫",
    "sleeping": "😴", "睡觉": "😴",
    "astonished": "😲", "惊讶": "😲",
    "flushed": "😳", "脸红": "😳",
    "pleading": "🥺", "请求": "🥺",
    "nauseated": "🤢", "恶心": "🤢",
    "sick": "🤒", "生病": "🤒",
    "cowboy": "🤠", "牛仔": "🤠",
    "confused": "😕", "困惑": "😕",
    "satisfied": "�满足",
    "rainbow": "🌈", "彩虹": "🌈",
    "sparkles": "✨", "闪光": "✨",
    "boom": "💥", "爆炸": "💥",
    "collision": "💥", "碰撞": "💥",
    "analyst": "🕵️", "侦探": "🕵️",
    "princess": "👸", "公主": "👸",
    "guardsman": "💂", "守卫": "💂",
    "dancer": "💃", "舞者": "💃",
    "police": "👮", "警察": "👮",
    "construction": "👷", "建筑": "👷",
    "bride": "👰", "新娘": "👰",
    "angel": "👼", "天使": "👼",
    "Santa": "🎅", "圣诞老人": "🎅",
    "fairy": "🧚", "仙女": "🧚",
    "vampire": "🧛", "吸血鬼": "🧛",
    "mermaid": "🧜", "美人鱼": "🧜",
    "elf": "🧝", "精灵": "🧝",
    "genie": "🧞", "精灵": "🧞",
    "zombie": "🧟", "僵尸": "🧟",
    "massage": "💆", "按摩": "💆",
    "haircut": "💇", "理发": "💇",
    "walking": "🚶", "走路": "🚶",
    "runner": "🏃", "跑步": "🏃",
    "dancer": "💃", "跳舞": "💃",
    "bow": "🙇", "鞠躬": "🙇",
    "couple": "👫", "情侣": "👫",
    "family": "👪", "家庭": "👪",
    "muscle": "💪", "肌肉": "💪", "力量": "💪",
    "point_left": "👈", "左边": "👈",
    "point_right": "👉", "右边": "👉",
    "point_up_2": "👆", "上边": "👆",
    "middle_finger": "🖕",
    "point_down": "👇", "下边": "👇",
    "v": "✌️", "胜利": "✌️",
    "hand_splayed": "🖐️",
    "metal": "🤘",
    "ok_hand": "👌", "完美": "👌",
    "raised_hand": "✋",
    "vulcan": "🖖",
    "writing": "✍️", "写作": "✍️",
    "pray": "🙏", "祈祷": "🙏",
    "handshake": "🤝", "握手": "🤝",
    "拳": "👊", "拳头": "👊",
    "facepunch": "🤜",
    "right_fist": "🤛",
    "oncoming_fist": "👊",
    "left_fist": "🤛",
}


def _add_emoji(text: str) -> str:
    """根据文本内容自动添加相关 emoji。
    
    策略：
    1. 先检查文本是否已有 emoji
    2. 查找关键词匹配
    3. 根据语义选择最合适的 emoji
    """
    import re as _re
    # 如果已有 emoji，不添加
    if _re.search(r'[\U0001F000-\U0001FFFF]|[\u2600-\u26FF]|[\u2700-\u27BF]', text):
        return text
    
    text_lower = text.lower()
    best_emoji = None
    best_score = 0
    
    for keyword, emoji in _EMOJI_KEYWORD_MAP.items():
        if keyword in text_lower:
            score = len(keyword)  # 越长越精确
            if score > best_score:
                best_score = score
                best_emoji = emoji
    
    if best_emoji:
        return f"{best_emoji}  {text}"
    return text


def _add_emoji_to_point(point: dict) -> dict:
    """为 heading 和 bullets 添加 emoji。"""
    result = dict(point)
    heading = result.get("heading", "")
    if heading:
        result["heading"] = _add_emoji(heading)
    bullets = result.get("bullets", [])
    if bullets:
        result["bullets"] = [_add_emoji(b) for b in bullets]
    return result


def _add_emoji_to_data(data: dict) -> dict:
    """为整个数据添加 emoji。"""
    result = dict(data)
    title = result.get("title", "")
    if title:
        result["title"] = _add_emoji(title)
    subtitle = result.get("subtitle", "")
    if subtitle:
        result["subtitle"] = _add_emoji(subtitle)
    summary = result.get("summary", "")
    if summary:
        result["summary"] = _add_emoji(summary)
    points = result.get("points", [])
    if points:
        result["points"] = [_add_emoji_to_point(p) for p in points]
    return result


def _unique_path(base: str) -> str:
    """🔴 2026-08-06 更名保留：路徑已存在 → 回傳 _v2/_v3... 版本化路徑（不覆蓋舊檔）。

    例：標題.pptx 已存在 → 標題_v2.pptx → 標題_v3.pptx…
    """
    if not os.path.exists(base):
        return base
    root, ext = os.path.splitext(base)
    i = 2
    while os.path.exists(f"{root}_v{i}{ext}"):
        i += 1
    return f"{root}_v{i}{ext}"


# ---------------------------------------------------------------------------
# Build PPT
# ---------------------------------------------------------------------------
def _add_title_slide(prs: Presentation, data: dict):
    """Add a dark-themed title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _DARK_BG

    # Title
    left, top, width, height = Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "Presentation")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    left, top, width, height = Inches(1.5), Inches(3.8), Inches(7), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("subtitle", "")
    p.font.size = Pt(24)
    p.font.color.rgb = _TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    left, top, width, height = Inches(3.5), Inches(3.5), Inches(3), Inches(0.05)
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()


def _add_content_slide(prs: Presentation, point: dict, index: int):
    """Add a content slide with heading and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT
    p.alignment = PP_ALIGN.LEFT

    # Heading
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # Accent line under heading
    left, top, width, height = Inches(0.6), Inches(2.3), Inches(2), Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()

    # Bullets
    bullets = point.get("bullets", [])
    left, top, width, height = Inches(0.8), Inches(2.7), Inches(8.4), Inches(2.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = _TEXT_LIGHT
        p.space_after = Pt(14)


def _add_hook_slide(prs: Presentation, point: dict):
    """hook 版型：大標題置中（吸睛開場），無編號 badge、無 bullets 列表。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # 大標題（置中）
    left, top, width, height = Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 補充（可選 1 條）
    bullets = point.get("bullets", [])
    if bullets:
        left, top, width, height = Inches(1.5), Inches(3.8), Inches(7), Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullets[0]
        p.font.size = Pt(22)
        p.font.color.rgb = _TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

    # Accent line
    left, top, width, height = Inches(3.5), Inches(3.5), Inches(3), Inches(0.05)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()


def _add_quote_slide(prs: Presentation, point: dict):
    """quote 版型：金句置中（大字形），下方來源/補充小字。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # 引號裝飾
    left, top, width, height = Inches(0.8), Inches(0.6), Inches(1.2), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.text = "「"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = _ACCENT

    # 金句（置中 36pt）
    left, top, width, height = Inches(1.2), Inches(1.6), Inches(7.6), Inches(2.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(48)

    # 補充/出處
    bullets = point.get("bullets", [])
    if bullets:
        left, top, width, height = Inches(2), Inches(4.0), Inches(6), Inches(0.7)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"— {bullets[0]}"
        p.font.size = Pt(20)
        p.font.color.rgb = _TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER


def _split_number(bullet: str) -> tuple[str | None, str]:
    """拆「70% 的使用者一週內回訪」→ ('70%', '的使用者一週內回訪')。

    找開頭數字+單位（%、倍、萬、億、人、台、元、x、×…）。無數字 → (None, 原樣)。
    """
    m = re.match(r'^([\d,]+\.?\d*\s*(?:%|％|倍|萬|億|人|台|元|x|×|次|個|張|件)?)', bullet)
    if m and m.group(1).strip():
        num = m.group(1).strip()
        rest = bullet[len(m.group(0)):].strip()
        return num, rest
    return None, bullet


def _add_data_slide(prs: Presentation, point: dict, index: int):
    """data 版型：編號 + 標題 + 數字卡（數字大字暖紅 + 說明小字）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT
    p.alignment = PP_ALIGN.LEFT

    # Heading
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # Accent line
    left, top, width, height = Inches(0.6), Inches(2.3), Inches(2), Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()

    # 數字卡（每 bullet 一張：數字大字 + 說明）— 🔴 2026-08-06 動態寬度：
    #   n=2 → 4.1 寬；n≥3 → 均分剩餘寬度（(8.8-(n-1)*0.3)/n），避免 3 卡超界
    bullets = point.get("bullets", [])
    n = max(len(bullets), 1)
    gap = Inches(0.3)
    usable = Inches(8.8)
    card_w = min(Inches(4.1), (usable - (n - 1) * gap) / n)
    total_w = n * card_w + (n - 1) * gap
    x0 = Inches(0.6) + (usable - total_w) / 2
    y0, card_h = Inches(2.8), Inches(2.2)
    for i, bullet in enumerate(bullets[:4]):
        num, desc = _split_number(bullet)
        cx = x0 + i * (card_w + gap)
        # 卡片背景
        card = slide.shapes.add_shape(1, cx, y0, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = _CARD_BG
        card.line.fill.background()
        # 數字
        txBox = slide.shapes.add_textbox(cx + Inches(0.25), y0 + Inches(0.25), card_w - Inches(0.5), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num if num else "—"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = _ACCENT if num else _TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        # 說明
        txBox = slide.shapes.add_textbox(cx + Inches(0.25), y0 + Inches(1.15), card_w - Inches(0.5), Inches(0.9))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc if desc else bullet
        p.font.size = Pt(20)
        p.font.color.rgb = _TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER


def _add_qa_slide(prs: Presentation, point: dict, index: int):
    """qa 版型：編號 + 標題 + 3 建議問題 + CTA。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT

    # Heading
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # 問題清單（前 3 條）＋ 最後一條 CTA
    bullets = point.get("bullets", [])
    questions = [b for b in bullets if not b.startswith(("→", "->", "現在", "立刻", "行動", "開始"))][:3]
    ctas = [b for b in bullets if b.startswith(("→", "->", "現在", "立刻", "行動", "開始"))]

    left, top, width, height = Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, q in enumerate(questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"❓  {q}"
        p.font.size = Pt(24)
        p.font.color.rgb = _TEXT_LIGHT
        p.space_after = Pt(12)

    # CTA（暖紅粗體）
    if ctas:
        left, top, width, height = Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctas[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = _ACCENT
        p.alignment = PP_ALIGN.CENTER


def _add_action_slide(prs: Presentation, point: dict, index: int):
    """action 版型：編號 + 標題 + 行動步驟（→ 前綴）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT

    # Heading
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # Accent line
    left, top, width, height = Inches(0.6), Inches(2.3), Inches(2), Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()

    # 行動步驟（→ 前綴，字大一點）
    bullets = point.get("bullets", [])
    left, top, width, height = Inches(0.8), Inches(2.7), Inches(8.4), Inches(2.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"→  {bullet}"
        p.font.size = Pt(26)
        p.font.color.rgb = _TEXT_LIGHT
        p.space_after = Pt(20)


def _add_summary_slide(prs: Presentation, data: dict):
    """Add a summary/conclusion slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _DARK_BG

    # "Key Takeaway" label
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY TAKEAWAY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = _ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    left, top, width, height = Inches(3.5), Inches(2.2), Inches(3), Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()

    # Summary text
    left, top, width, height = Inches(1), Inches(2.6), Inches(8), Inches(3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("summary", "")
    p.font.size = Pt(26)
    p.font.color.rgb = _TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(36)


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------
def generate_ppt(script: str, title: str, lang: str = "zh", out_dir: str = ".", scheme: str = "dark") -> str | None:
    """Generate a PowerPoint presentation from a podcast script.
    
    Args:
        script: podcast script text
        title: video title
        lang: target language
        out_dir: output directory
    
    Returns:
        Path to the generated .pptx file, or None on failure.
    """
    print("[INFO] Extracting key points for PPT...", file=sys.stderr)
    data = _extract_key_points(script, title, lang)
    if data:
        data = _add_emoji_to_data(data)
    if not data:
        print("[WARN] Could not extract key points, using fallback", file=sys.stderr)
        data = {
            "title": title,
            "subtitle": "",
            "points": [{"heading": "重點整理", "bullets": ["無法提取結構化重點"]}],
            "summary": "請參閱口播腳本了解完整內容。",
        }

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Build slides — 🔴 2026-08-06 階段 2：依 slide_type 選版型（hook/quote/data/qa/action/content）
    _add_title_slide(prs, data)
    for i, point in enumerate(data.get("points", []), 1):
        st = point.get("slide_type", "content")
        if st == "hook":
            _add_hook_slide(prs, point)
        elif st == "quote":
            _add_quote_slide(prs, point)
        elif st == "data":
            _add_data_slide(prs, point, i)
        elif st == "qa":
            _add_qa_slide(prs, point, i)
        elif st == "action":
            _add_action_slide(prs, point, i)
        elif st == "comparison":
            _add_comparison_slide(prs, point, i)
        elif st == "timeline":
            _add_timeline_slide(prs, point, i)
        elif st == "split":
            _add_split_slide(prs, point, i)
        else:
            _add_content_slide(prs, point, i)
    _add_summary_slide(prs, data)

    # Apply color scheme
    scheme_colors = COLOR_SCHEMES.get(scheme, COLOR_SCHEMES["dark"])
    # 更新 background
    for slide in prs.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = scheme_colors["bg"]
    # 更新文字色
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        c = run.font.color.rgb
                        if c == RGBColor(0xFF, 0xFF, 0xFF):
                            run.font.color.rgb = scheme_colors["text_white"]
                        elif c == RGBColor(0xCC, 0xCC, 0xCC):
                            run.font.color.rgb = scheme_colors["text_light"]
                        elif c == RGBColor(0xE8, 0x4D, 0x3D):
                            run.font.color.rgb = scheme_colors["accent"]
                    except AttributeError:
                        pass
    
    # Save
    safe_title = title.replace("/", "_").replace("\\", "_")[:80]
    # 🔴 2026-08-06：繁中字型（Noto Sans CJK TC + a:ea 属性）
    _apply_cjk_font(prs)
    # 🔴 2026-08-06：更名保留 — 同名已存在 → _v2/_v3（不覆蓋舊檔）
    pptx_path = _unique_path(os.path.join(out_dir, f"{safe_title}.pptx"))
    prs.save(pptx_path)
    os.chmod(pptx_path, 0o777)
    print(f"[OK] PPT saved: {pptx_path} ({len(prs.slides)} slides)", file=sys.stderr)
    return pptx_path


def _add_comparison_slide(prs: Presentation, point: dict, index: int):
    """comparison 版型：左右對比（A vs B）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT

    # Heading
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # Accent line
    left, top, width, height = Inches(0.6), Inches(2.3), Inches(2), Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()

    # 左右分欄
    bullets = point.get("bullets", [])
    left_col, right_col = Inches(0.6), Inches(5.4)
    width_col = Inches(4.2)
    y0, height = Inches(2.7), Inches(2.5)

    for i, bullet in enumerate(bullets[:4]):
        if "→" in bullet:
            left_text, _, right_text = bullet.partition("→")
            # 左欄
            txBox = slide.shapes.add_textbox(left_col, y0 + i * Inches(0.65), width_col, Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"❶ {left_text.strip()}"
            p.font.size = Pt(22)
            p.font.color.rgb = _TEXT_LIGHT
            # 右欄
            txBox = slide.shapes.add_textbox(right_col, y0 + i * Inches(0.65), width_col, Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"❷ {right_text.strip()}"
            p.font.size = Pt(22)
            p.font.color.rgb = _ACCENT
        else:
            # 沒有分隔 → 統一放左欄
            txBox = slide.shapes.add_textbox(left_col, y0 + i * Inches(0.65), width_col * 2, Inches(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = _TEXT_LIGHT


def _add_timeline_slide(prs: Presentation, point: dict, index: int):
    """timeline 版型：時間軸/流程（箭頭連結）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT

    # Heading
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # Accent line
    left, top, width, height = Inches(0.6), Inches(2.3), Inches(2), Inches(0.04)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _ACCENT
    shape.line.fill.background()

    # 流程箭頭
    bullets = point.get("bullets", [])
    left, top, width, height = Inches(0.8), Inches(2.7), Inches(8.4), Inches(2.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        arrow = "➜" if i > 0 else "①"
        p.text = f"{arrow}  {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = _TEXT_LIGHT
        p.space_after = Pt(14)


def _add_split_slide(prs: Presentation, point: dict, index: int):
    """split 版型：左標題 + 右內容（案例/詳情）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _DARK_BG

    # Number badge
    left, top, width, height = Inches(0.6), Inches(0.4), Inches(0.8), Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{index:02d}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _ACCENT

    # 左區大標題
    left, top, width, height = Inches(0.6), Inches(1.4), Inches(3.8), Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = point.get("heading", "")
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = _TEXT_WHITE

    # 右區內容
    bullets = point.get("bullets", [])
    left, top, width, height = Inches(5.0), Inches(1.4), Inches(4.4), Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = _TEXT_LIGHT
        p.space_after = Pt(12)
