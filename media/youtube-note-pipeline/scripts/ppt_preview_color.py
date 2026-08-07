#!/usr/bin/env python3
"""彩色 PPT 預覽渲染器：從 PPTX 抽文字 → 生成 HTML → Playwright 截圖

為什麼需要這個：Pillow 無法渲染彩色 emoji 字型（NotoColorEmoji.ttf 是 CBDT
bitmap 格式 → `ImageFont.truetype` 報 "invalid pixel size"；NotoEmoji-Regular.ttf
是單色白線條）。Playwright + Chromium 能完整渲染彩色 emoji（💊 黃紅膠囊、🔬
青藍顯微鏡），所以用 HTML 中繼 + Chromium 截圖做出「後製套色」的預覽圖。

用途：交付 .pptx 給使用者前，先跑這個產出彩色 PNG 預覽（Telegram 對 .pptx
無視覺預覽，見 pitfall 60）。

需求：node + playwright（本機在 /opt/data/tmp/node_modules/playwright）。
用法：
    python ppt_preview_color.py <file.pptx> <out_dir> [dark|blue|green|light]
"""
import os
import sys
import subprocess

sys.path.insert(0, "/opt/data/skills/media/youtube-note-pipeline/scripts")
from pptx import Presentation  # noqa: E402


def extract_slides(pptx_path):
    """從 PPTX 讀出每頁文字（text、size、bold、color）。"""
    prs = Presentation(pptx_path)
    slides = []
    for slide in prs.slides:
        items = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(r.text for r in para.runs)
                if not txt.strip():
                    continue
                size = None
                bold = False
                color = None
                for r in para.runs:
                    if r.font.size:
                        size = r.font.size.pt
                    bold = bold or bool(r.font.bold)
                    try:
                        if r.font.color and r.font.color.type is not None:
                            color = str(r.font.color.rgb)
                    except Exception:
                        pass
                items.append({
                    "text": txt.strip(),
                    "size": size or 18,
                    "bold": bold,
                    "color": color,
                })
        slides.append(items)
    return slides


def slide_to_html(slide_items, slide_no, total, bg="#1a1a2e", accent="#E84D3D"):
    parts = [f"<div class='slide' id='s{slide_no}'>"]
    parts.append(f"<div class='page-no'>{slide_no:02d} / {total}</div>")
    for item in slide_items:
        size = item["size"]
        cls = "title" if (item["bold"] and size >= 24) else ("bold" if item["bold"] else "body")
        style = f"font-size:{size * 1.25:.0f}px;"
        if item["color"]:
            style += f"color:#{item['color']};"
        parts.append(f"<div class='{cls}' style='{style}'>{item['text']}</div>")
    parts.append("</div>")
    return "\n".join(parts)


def build_html(slides, bg="#1a1a2e", accent="#E84D3D"):
    css = f"""
    body {{ margin:0; background:{bg}; font-family:'Noto Sans TC','Source Han Sans TC',sans-serif; }}
    .slide {{ width:1280px; height:720px; background:{bg}; color:#fff; padding:60px 80px;
             box-sizing:border-box; display:flex; flex-direction:column; gap:14px;
             justify-content:center; page-break-after:always; position:relative; }}
    .page-no {{ position:absolute; top:20px; right:40px; font-size:16px; color:rgba(255,255,255,.4); }}
    .title {{ font-size:44px; font-weight:700; color:{accent}; }}
    .bold {{ font-size:30px; font-weight:700; color:#fff; }}
    .body {{ font-size:26px; color:#ccc; }}
    """
    body = "\n".join(slide_to_html(s, i, len(slides), bg, accent)
                     for i, s in enumerate(slides, 1))
    return (f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
            f"<style>{css}</style></head><body>{body}</body></html>")


def render_playwright(html_path, out_dir, total):
    """用 node playwright 截圖（每頁一張）。.slide 需 position:relative，
    否則 .page-no 的 absolute 定位會跑到整份 HTML 的角落。"""
    node_script = f"""
    const {{ chromium }} = require('/opt/data/tmp/node_modules/playwright');
    (async () => {{
        const b = await chromium.launch();
        const page = await b.newPage({{ viewport: {{ width: 1280, height: 720 }} }});
        await page.goto('file://{html_path}');
        await page.waitForTimeout(300);
        for (let i = 1; i <= {total}; i++) {{
            const el = await page.$('#s' + i);
            if (el) {{ await el.screenshot({{ path: '{out_dir}/slide_' + String(i).padStart(2,'0') + '.png' }}); }}
        }}
        await b.close();
        console.log('DONE');
    }})();
    """
    r = subprocess.run(["node", "-e", node_script],
                       capture_output=True, text=True, timeout=120)
    if r.returncode != 0:
        print("playwright 失敗:", r.stderr[-500:], file=sys.stderr)
        return False
    return True


def main():
    pptx = sys.argv[1]
    out_dir = sys.argv[2]
    scheme = sys.argv[3] if len(sys.argv) > 3 else "dark"
    os.makedirs(out_dir, exist_ok=True)

    schemes = {
        "dark":  ("#1a1a2e", "#E84D3D"),
        "blue":  ("#0f1c3f", "#4D9BE8"),
        "green": ("#0d2b1e", "#4DC878"),
        "light": ("#f5f5f5", "#E84D3D"),
    }
    bg, accent = schemes.get(scheme, schemes["dark"])

    slides = extract_slides(pptx)
    html = build_html(slides, bg, accent)
    html_path = os.path.join(out_dir, "preview.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"[INFO] HTML 生成: {html_path} ({len(slides)} 頁)", file=sys.stderr)

    if render_playwright(html_path, out_dir, len(slides)):
        print(f"[OK] 彩色預覽圖完成: {out_dir}", file=sys.stderr)
    else:
        sys.exit(1)


if __name__ == "__main__":
    main()
